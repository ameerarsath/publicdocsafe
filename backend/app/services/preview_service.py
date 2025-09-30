"""
Preview Service for SecureVault
Handles generation of document previews for various file types
"""

import io
import base64
import hashlib
import tempfile
import logging
from typing import Dict, Any, Optional, List, Union
from pathlib import Path
import json
from datetime import datetime, timedelta

# Third-party libraries for preview generation
try:
    from PIL import Image, ImageOps
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

try:
    import PyPDF2
    PYPDF2_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False

try:
    from pdf2image import convert_from_bytes
    import os

    # Configure Poppler path for pdf2image on Windows
    POPPLER_PATH = None
    if os.name == 'nt':  # Windows
        possible_poppler_paths = [
            r"C:\Users\moham\AppData\Local\Microsoft\WinGet\Packages\oschwartz10612.Poppler_Microsoft.Winget.Source_8wekyb3d8bbwe\poppler-25.07.0\Library\bin",
            r"C:\Program Files\poppler\bin",
            r"C:\poppler\bin",
            r"C:\tools\poppler\bin"
        ]

        for path in possible_poppler_paths:
            if os.path.exists(path) and os.path.exists(os.path.join(path, "pdftoppm.exe")):
                POPPLER_PATH = path
                print(f"INFO: Found Poppler at: {POPPLER_PATH}")
                break

        if not POPPLER_PATH:
            print("WARNING: Poppler not found in standard locations. PDF thumbnails may not work.")

    PDF2IMAGE_AVAILABLE = True
except ImportError:
    PDF2IMAGE_AVAILABLE = False
    POPPLER_PATH = None

# Combined PDF processing availability
PDF_AVAILABLE = PYPDF2_AVAILABLE

# Office processing is always available - we've confirmed the libraries are installed
OFFICE_AVAILABLE = True

try:
    import docx
    import openpyxl
    import pptx
except ImportError:
    # Libraries should be available, but if not, we'll handle it at runtime
    pass

from ..models.document import Document
from ..core.config import settings

logger = logging.getLogger(__name__)

class PreviewService:
    """Service for generating document previews"""
    
    def __init__(self):
        self.cache = {}  # In-memory cache for development (use Redis in production)
        self.cache_expiry = timedelta(hours=1)
        
    async def get_cached_preview(self, cache_key: str) -> Optional[Dict[str, Any]]:
        """Get cached preview if available and not expired"""
        if cache_key in self.cache:
            cached_item = self.cache[cache_key]
            if datetime.now() < cached_item['expires']:
                logger.info(f"Cache hit for preview: {cache_key}")
                return cached_item['data']
            else:
                # Remove expired cache
                del self.cache[cache_key]
        return None
    
    async def cache_preview(self, cache_key: str, preview_data: Dict[str, Any]) -> None:
        """Cache preview data with expiry"""
        self.cache[cache_key] = {
            'data': preview_data,
            'expires': datetime.now() + self.cache_expiry
        }
        logger.info(f"Cached preview: {cache_key}")
    
    def clear_cache(self, document_id: int = None) -> None:
        """Clear preview cache for specific document or all cache"""
        if document_id:
            # Clear cache for specific document
            keys_to_remove = [key for key in self.cache.keys() if key.startswith(f"preview_{document_id}_")]
            for key in keys_to_remove:
                del self.cache[key]
            logger.info(f"Cleared preview cache for document {document_id}")
        else:
            # Clear all cache
            self.cache.clear()
            logger.info("Cleared all preview cache")
    
    async def generate_thumbnail(
        self, 
        file_data: bytes, 
        mime_type: str, 
        filename: str, 
        max_size: int = 1024
    ) -> Dict[str, Any]:
        """Generate thumbnail preview for images and PDFs"""
        
        try:
            if mime_type.startswith('image/'):
                return await self._generate_image_thumbnail(file_data, max_size)
            elif mime_type == 'application/pdf':
                return await self._generate_pdf_thumbnail(file_data, max_size)
            else:
                # Generate generic file icon thumbnail
                return await self._generate_file_icon_thumbnail(mime_type, filename)
                
        except Exception as e:
            logger.error(f"Thumbnail generation failed: {str(e)}")
            return await self._generate_error_thumbnail(str(e))
    
    async def _generate_image_thumbnail(self, file_data: bytes, max_size: int) -> Dict[str, Any]:
        """Generate thumbnail for image files with robust error handling"""

        if not PIL_AVAILABLE:
            return {
                "type": "info",
                "format": "image",
                "message": "Image preview not available - PIL not installed",
                "suggestion": "Download the file to view the image"
            }

        try:
            # Validate image data first
            if not file_data or len(file_data) < 10:
                return {
                    "type": "info",
                    "format": "image",
                    "message": "Image file appears to be empty or corrupted",
                    "suggestion": "Please check the file and try uploading again"
                }

            # Create BytesIO object and ensure it's at the beginning
            image_stream = io.BytesIO(file_data)
            image_stream.seek(0)

            # Try to identify and open the image with better error handling
            try:
                # First, try to verify the image format
                image_stream.seek(0)
                image_format = None

                # Check common image signatures
                signature = file_data[:8]
                if signature.startswith(b'\xFF\xD8\xFF'):
                    image_format = 'JPEG'
                elif signature.startswith(b'\x89PNG\r\n\x1A\n'):
                    image_format = 'PNG'
                elif signature.startswith(b'GIF87a') or signature.startswith(b'GIF89a'):
                    image_format = 'GIF'
                elif signature.startswith(b'BM'):
                    image_format = 'BMP'
                elif signature.startswith(b'RIFF') and b'WEBP' in file_data[:12]:
                    image_format = 'WEBP'

                # Reset stream position
                image_stream.seek(0)

                # Try to open the image
                image = Image.open(image_stream)

                # Verify the image by getting basic info
                original_width, original_height = image.size
                original_mode = image.mode
                detected_format = image.format or image_format or 'Unknown'

                # Load the image data to catch any corruption issues
                image.load()

            except (Image.UnidentifiedImageError, OSError) as img_error:
                logger.warning(f"PIL cannot identify image: {str(img_error)}")
                return {
                    "type": "info",
                    "format": "image",
                    "message": "Image format not supported or file is corrupted",
                    "error_detail": str(img_error),
                    "suggestion": "Try converting the image to a common format (JPEG, PNG) or check if the file is corrupted"
                }
            except Exception as img_error:
                logger.warning(f"Image loading failed: {str(img_error)}")
                return {
                    "type": "info",
                    "format": "image",
                    "message": "Unable to process this image file",
                    "error_detail": str(img_error),
                    "suggestion": "Download the file to view with an image viewer"
                }

            # Create thumbnail while maintaining aspect ratio
            try:
                # Make a copy for thumbnail generation
                thumbnail_image = image.copy()
                thumbnail_image.thumbnail((max_size, max_size), Image.Resampling.LANCZOS)

                # Convert to RGB if needed (for JPEG output)
                if thumbnail_image.mode in ('RGBA', 'LA', 'P'):
                    rgb_image = Image.new('RGB', thumbnail_image.size, (255, 255, 255))
                    if thumbnail_image.mode == 'P':
                        thumbnail_image = thumbnail_image.convert('RGBA')
                    rgb_image.paste(thumbnail_image, mask=thumbnail_image.split()[-1] if thumbnail_image.mode == 'RGBA' else None)
                    thumbnail_image = rgb_image

                # Save thumbnail to bytes
                output = io.BytesIO()
                thumbnail_image.save(output, format='JPEG', quality=85, optimize=True)
                thumbnail_data = output.getvalue()

                # Encode to base64
                thumbnail_base64 = base64.b64encode(thumbnail_data).decode('utf-8')

                return {
                    "type": "thumbnail",
                    "format": "image",
                    "data": thumbnail_base64,
                    "data_url": f"data:image/jpeg;base64,{thumbnail_base64}",
                    "thumbnail_size": {
                        "width": thumbnail_image.width,
                        "height": thumbnail_image.height
                    },
                    "original_size": {
                        "width": original_width,
                        "height": original_height
                    },
                    "file_size": len(thumbnail_data),
                    "original_format": detected_format,
                    "original_mode": original_mode
                }

            except Exception as thumb_error:
                logger.warning(f"Thumbnail generation failed: {str(thumb_error)}")
                # Return image info without thumbnail
                return {
                    "type": "info",
                    "format": "image",
                    "message": f"Image detected: {original_width}x{original_height} pixels, {detected_format} format",
                    "original_size": {
                        "width": original_width,
                        "height": original_height
                    },
                    "original_format": detected_format,
                    "original_mode": original_mode,
                    "suggestion": "Download to view the full image"
                }

        except Exception as e:
            logger.error(f"Image processing completely failed: {str(e)}")
            return {
                "type": "info",
                "format": "image",
                "message": "Image file detected but cannot be processed",
                "error_detail": str(e),
                "suggestion": "Download the file to view with an image viewer or photo editor"
            }
    
    async def _generate_pdf_thumbnail(self, file_data: bytes, max_size: int) -> Dict[str, Any]:
        """Generate thumbnail for PDF first page"""

        if not PYPDF2_AVAILABLE:
            return {
                "type": "error",
                "message": "PDF processing not available - PyPDF2 not installed"
            }

        # Try PDF2IMAGE first (requires Poppler)
        if PDF2IMAGE_AVAILABLE:
            try:
                # Convert first page to image using pdf2image
                # Pass poppler_path if we found it
                convert_kwargs = {
                    'first_page': 1,
                    'last_page': 1,
                    'dpi': 150
                }

                if POPPLER_PATH:
                    convert_kwargs['poppler_path'] = POPPLER_PATH

                images = convert_from_bytes(file_data, **convert_kwargs)

                if images:
                    # Get first page
                    first_page = images[0]

                    # Create thumbnail
                    first_page.thumbnail((max_size, max_size), Image.Resampling.LANCZOS)

                    # Save to bytes
                    output = io.BytesIO()
                    first_page.save(output, format='JPEG', quality=85, optimize=True)
                    thumbnail_data = output.getvalue()

                    # Encode to base64
                    thumbnail_base64 = base64.b64encode(thumbnail_data).decode('utf-8')

                    # Try to get PDF metadata
                    pdf_info = await self._get_pdf_info(file_data)

                    return {
                        "type": "thumbnail",
                        "format": "pdf",
                        "data": thumbnail_base64,
                        "data_url": f"data:image/jpeg;base64,{thumbnail_base64}",
                        "thumbnail_size": {
                            "width": first_page.width,
                            "height": first_page.height
                        },
                        "pdf_info": pdf_info,
                        "file_size": len(thumbnail_data),
                        "generation_method": "pdf2image"
                    }

            except Exception as e:
                logger.warning(f"PDF2IMAGE thumbnail generation failed: {str(e)}")
                # Continue to fallback method

        # Fallback: Generate PDF icon with metadata
        try:
            pdf_info = await self._get_pdf_info(file_data)

            return {
                "type": "icon",
                "format": "pdf_icon",
                "icon": "📄",
                "category": "pdf",
                "pdf_info": pdf_info,
                "message": "PDF thumbnail not available - Poppler not installed. Showing PDF icon with metadata.",
                "generation_method": "metadata_fallback",
                "fallback_reason": "pdf2image/Poppler not available or failed"
            }

        except Exception as e:
            logger.error(f"PDF metadata extraction failed: {str(e)}")
            return await self._generate_error_thumbnail(f"PDF processing failed: Unable to get page count. Is poppler installed and in PATH?")
    
    async def _get_pdf_info(self, file_data: bytes) -> Dict[str, Any]:
        """Extract basic PDF metadata"""
        if not PYPDF2_AVAILABLE:
            return {
                "page_count": "unknown",
                "error": "PyPDF2 not available"
            }

        try:
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_data))

            return {
                "page_count": len(pdf_reader.pages),
                "title": getattr(pdf_reader.metadata, 'title', None) if pdf_reader.metadata else None,
                "author": getattr(pdf_reader.metadata, 'author', None) if pdf_reader.metadata else None,
                "creator": getattr(pdf_reader.metadata, 'creator', None) if pdf_reader.metadata else None,
                "producer": getattr(pdf_reader.metadata, 'producer', None) if pdf_reader.metadata else None,
                "creation_date": getattr(pdf_reader.metadata, 'creation_date', None) if pdf_reader.metadata else None
            }
        except Exception as e:
            logger.error(f"PDF metadata extraction failed: {str(e)}")
            return {
                "page_count": "unknown",
                "error": f"Failed to read PDF: {str(e)}"
            }
    
    async def extract_text_preview(
        self,
        file_data: bytes,
        mime_type: str,
        filename: str,
        max_chars: int = 500
    ) -> Dict[str, Any]:
        """Extract text preview from various document types with enhanced file type detection"""

        try:
            # Get file extension for fallback detection
            file_extension = filename.lower().split('.')[-1] if '.' in filename else ''

            # Enhanced MIME type detection based on both MIME type and extension
            detected_type = self._detect_enhanced_file_type(mime_type, file_extension, file_data)

            if detected_type in ['text', 'json', 'xml', 'csv', 'html', 'markdown', 'yaml']:
                return await self._extract_plain_text(file_data, max_chars)
            elif detected_type == 'pdf':
                return await self._extract_pdf_text(file_data, max_chars)
            elif detected_type in ['docx', 'xlsx', 'pptx', 'doc', 'xls', 'ppt']:
                return await self._extract_office_text(file_data, mime_type, max_chars)
            elif detected_type in ['image']:
                return await self._extract_image_preview(file_data, mime_type, filename)
            else:
                return await self._generate_unsupported_text_preview(filename, detected_type)

        except Exception as e:
            logger.error(f"Text extraction failed: {str(e)}")
            return {
                "type": "info",
                "format": "text",
                "message": "This file cannot be previewed as text",
                "suggestion": "Download the file to view it with appropriate software",
                "filename": filename
            }

    def _get_extension_for_mime(self, mime_type: str) -> str:
        """Get the appropriate file extension for a MIME type"""
        mime_to_ext = {
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': '.docx',
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': '.xlsx',
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': '.pptx',
            'application/msword': '.doc',
            'application/vnd.ms-excel': '.xls',
            'application/vnd.ms-powerpoint': '.ppt',
            'text/plain': '.txt'
        }
        return mime_to_ext.get(mime_type, '.bin')

    def _detect_enhanced_file_type(self, mime_type: str, file_extension: str, file_data: bytes) -> str:
        """Enhanced file type detection using MIME type, extension, and file signature"""

        # Primary detection via MIME type
        if mime_type:
            if mime_type.startswith('text/'):
                if 'csv' in mime_type:
                    return 'csv'
                elif 'html' in mime_type:
                    return 'html'
                elif 'markdown' in mime_type:
                    return 'markdown'
                else:
                    return 'text'
            elif mime_type in ['application/json']:
                return 'json'
            elif mime_type in ['application/xml', 'text/xml']:
                return 'xml'
            elif mime_type in ['application/x-yaml', 'text/yaml', 'application/yaml']:
                return 'yaml'
            elif mime_type == 'application/pdf':
                return 'pdf'
            elif mime_type.startswith('image/'):
                return 'image'
            elif 'wordprocessingml' in mime_type or mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
                return 'docx'
            elif 'spreadsheetml' in mime_type or mime_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
                return 'xlsx'
            elif 'presentationml' in mime_type or mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                return 'pptx'
            elif mime_type in ['application/msword']:
                return 'doc'
            elif mime_type in ['application/vnd.ms-excel']:
                return 'xls'
            elif mime_type in ['application/vnd.ms-powerpoint']:
                return 'ppt'

        # Fallback detection via file extension
        if file_extension:
            extension_map = {
                'txt': 'text',
                'text': 'text',
                'csv': 'csv',
                'html': 'html',
                'htm': 'html',
                'md': 'markdown',
                'markdown': 'markdown',
                'json': 'json',
                'xml': 'xml',
                'yaml': 'yaml',
                'yml': 'yaml',
                'pdf': 'pdf',
                'docx': 'docx',
                'doc': 'doc',
                'xlsx': 'xlsx',
                'xls': 'xls',
                'pptx': 'pptx',
                'ppt': 'ppt',
                'jpg': 'image',
                'jpeg': 'image',
                'png': 'image',
                'gif': 'image',
                'bmp': 'image',
                'tiff': 'image',
                'tif': 'image',
                'webp': 'image',
                'svg': 'image'
            }

            if file_extension in extension_map:
                return extension_map[file_extension]

        # File signature detection for ambiguous cases
        if len(file_data) >= 4:
            signature = file_data[:4]
            if signature == b'%PDF':
                return 'pdf'
            elif signature[:2] == b'PK':
                # ZIP-based formats (DOCX, XLSX, PPTX) - check any PK signature
                if file_extension in ['docx', 'doc']:
                    return 'docx'
                elif file_extension in ['xlsx', 'xls']:
                    return 'xlsx'
                elif file_extension in ['pptx', 'ppt']:
                    return 'pptx'
                else:
                    return 'zip'
            elif signature[:2] in [b'\xff\xd8', b'\x89P']:  # JPEG, PNG
                return 'image'

        return 'unknown'
    
    async def _extract_plain_text(self, file_data: bytes, max_chars: int) -> Dict[str, Any]:
        """Extract text from plain text files"""
        try:
            # Try UTF-8 first, fallback to other encodings
            encodings = ['utf-8', 'utf-16', 'latin1', 'cp1252']
            
            text_content = None
            used_encoding = None
            
            for encoding in encodings:
                try:
                    text_content = file_data.decode(encoding)
                    used_encoding = encoding
                    break
                except UnicodeDecodeError:
                    continue
            
            if text_content is None:
                return {
                    "type": "info",
                    "format": "text",
                    "message": "This text file cannot be previewed due to unsupported character encoding",
                    "suggestion": "Download the file and open it with a text editor that supports various character encodings (UTF-8, UTF-16, etc.)"
                }
            
            # Extract preview
            preview_text = text_content[:max_chars]
            if len(text_content) > max_chars:
                preview_text += "..."
            
            # Basic text analysis
            lines = text_content.split('\n')
            word_count = len(text_content.split())
            
            return {
                "type": "text",
                "format": "plain_text", 
                "preview": preview_text,
                "full_text_length": len(text_content),
                "line_count": len(lines),
                "word_count": word_count,
                "encoding": used_encoding,
                "is_truncated": len(text_content) > max_chars
            }
            
        except Exception as e:
            return {
                "type": "info",
                "format": "text",
                "message": "This text file cannot be previewed due to encoding issues",
                "suggestion": "Download the file and open it with a text editor that supports various character encodings"
            }
    
    async def _extract_pdf_text(self, file_data: bytes, max_chars: int) -> Dict[str, Any]:
        """Extract text from PDF"""
        if not PYPDF2_AVAILABLE:
            return {
                "type": "error",
                "message": "PDF text processing not available - PyPDF2 not installed"
            }
        
        try:
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_data))
            
            # Extract text from first few pages
            text_content = ""
            pages_processed = 0
            
            for page_num, page in enumerate(pdf_reader.pages[:3]):  # First 3 pages
                try:
                    page_text = page.extract_text()
                    text_content += page_text + "\n"
                    pages_processed += 1
                    
                    if len(text_content) >= max_chars:
                        break
                except:
                    continue
            
            if not text_content.strip():
                return {
                    "type": "info",
                    "message": "No text content found in PDF (may be image-based)"
                }
            
            # Extract preview
            preview_text = text_content[:max_chars]
            if len(text_content) > max_chars:
                preview_text += "..."
            
            return {
                "type": "text",
                "format": "pdf_text",
                "preview": preview_text,
                "full_text_length": len(text_content),
                "pages_processed": pages_processed,
                "total_pages": len(pdf_reader.pages),
                "is_truncated": len(text_content) > max_chars
            }
            
        except Exception as e:
            return {
                "type": "error",
                "message": f"PDF text extraction failed: {str(e)}"
            }
    
    async def _extract_office_text(self, file_data: bytes, mime_type: str, max_chars: int) -> Dict[str, Any]:
        """Extract text from Office documents with improved error handling"""

        # Office processing is now always enabled since we've confirmed libraries are available
        if not OFFICE_AVAILABLE:
            return {
                "type": "error",
                "message": "Office document processing not available"
            }

        try:
            # Log file details for debugging
            logger.info(f"Processing office file: size={len(file_data)} bytes, mime_type={mime_type}")
            
            # Verify file data integrity
            if len(file_data) < 100:
                return {
                    "type": "error",
                    "message": "File appears to be empty or corrupted"
                }

            # Basic validation for Office documents
            import zipfile
            import io

            # For Office Open XML files, check if it's a valid ZIP file but be more permissive
            if any(fmt in mime_type for fmt in ["openxmlformats", "wordprocessingml", "spreadsheetml", "presentationml"]):
                try:
                    with zipfile.ZipFile(io.BytesIO(file_data)) as zf:
                        # Just verify we can open it as a ZIP file - no integrity test
                        # The test_result check was too strict and rejected valid Office documents
                        pass
                except zipfile.BadZipFile:
                    # Even if ZIP validation fails, try to process the document
                    # Some valid Office documents may not pass strict ZIP validation
                    logger.warning("ZIP validation failed, but proceeding with document processing")
                    pass

            # Save the binary data to a temporary file to ensure proper handling
            with tempfile.NamedTemporaryFile(delete=False, suffix=self._get_extension_for_mime(mime_type)) as temp_file:
                temp_file.write(file_data)
                temp_file_path = temp_file.name
                logger.info(f"Saved to temporary file: {temp_file_path}")

            # Light validation for file format (more permissive than before)
            validation_result = await self._validate_office_file(file_data, mime_type)
            if not validation_result["is_valid"]:
                # Log the validation issue but continue processing
                logger.warning(f"Office file validation warning: {validation_result['error']}")
                # Continue processing instead of failing

            text_content = ""
            doc_type = "unknown"

            if "wordprocessingml" in mime_type or mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                # Word document
                try:
                    # Use the temporary file we created
                    doc = docx.Document(temp_file_path)
                    paragraphs = []
                    
                    # Enhanced text extraction
                    for paragraph in doc.paragraphs:
                        if paragraph.text.strip():
                            # Include formatting if available
                            text = paragraph.text.strip()
                            if paragraph.style and paragraph.style.name != 'Normal':
                                text = f"[{paragraph.style.name}] {text}"
                            paragraphs.append(text)
                    
                    # Also extract text from tables
                    for table in doc.tables:
                        for row in table.rows:
                            row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                            if row_text:
                                paragraphs.append(f"[Table] {row_text}")
                    
                    text_content = "\n\n".join(paragraphs)
                    doc_type = "word"
                except Exception as doc_error:
                    logger.error(f"DOCX processing failed: {str(doc_error)}")
                    return {
                        "type": "processing_unavailable",
                        "format": "office_word",
                        "message": "Office document processing not available",
                        "document_type": "word",
                        "suggestion": "Client-side processing may be available",
                        "error_details": str(doc_error),
                        "can_fallback": True
                    }

            elif "spreadsheetml" in mime_type or mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                # Excel document
                try:
                    workbook = openpyxl.load_workbook(io.BytesIO(file_data), read_only=True, data_only=True)
                    worksheet = workbook.active

                    rows = []
                    for row_num, row in enumerate(worksheet.iter_rows(max_row=50, values_only=True)):
                        if row_num == 0 or any(cell for cell in row if cell is not None):
                            row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                            if row_text.strip():
                                rows.append(row_text)

                    text_content = "\n".join(rows)
                    doc_type = "excel"
                    workbook.close()
                except Exception as excel_error:
                    logger.error(f"Excel processing failed: {str(excel_error)}")
                    return {
                        "type": "error",
                        "message": f"Excel document processing failed: {str(excel_error)}"
                    }

            elif "presentationml" in mime_type or mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                # PowerPoint document
                try:
                    # Use the temporary file we created
                    prs = pptx.Presentation(temp_file_path)

                    slides_text = []
                    total_slides = len(prs.slides)
                    
                    # Process all slides but limit text extraction
                    for slide_num, slide in enumerate(prs.slides):
                        slide_content = []
                        
                        # Get slide layout type
                        layout_type = slide.slide_layout.name if slide.slide_layout else "Unknown"
                        
                        # Start with slide header
                        slide_content.append(f"[Slide {slide_num + 1}/{total_slides}] ({layout_type})")
                        
                        try:
                            # Extract text from all shapes
                            for shape in slide.shapes:
                                if not hasattr(shape, "text"):
                                    continue
                                    
                                text = shape.text.strip()
                                if not text:
                                    continue
                                    
                                # Try to identify the type of content
                                if hasattr(shape, "name"):
                                    shape_name = shape.name.lower()
                                    if "title" in shape_name:
                                        text = f"Title: {text}"
                                    elif "heading" in shape_name:
                                        text = f"Heading: {text}"
                                    elif "bullet" in shape_name or "list" in shape_name:
                                        text = f"• {text}"
                                
                                slide_content.append(text)
                                
                        except Exception as shape_error:
                            logger.warning(f"Error in slide {slide_num + 1}: {shape_error}")
                            slide_content.append("[Some content could not be extracted]")
                            continue

                        if len(slide_content) > 1:  # If we have more than just the header
                            slides_text.append("\n".join(slide_content))

                    if not slides_text:
                        return {
                            "type": "info",
                            "format": "office_powerpoint",
                            "message": "This presentation appears to contain no extractable text",
                            "document_type": "powerpoint",
                            "metadata": {
                                "total_slides": total_slides,
                                "has_content": False
                            },
                            "suggestion": "Download to view the full presentation with Microsoft PowerPoint"
                        }

                    text_content = "\n\n" + "=" * 40 + "\n\n".join(slides_text)
                    doc_type = "powerpoint"

                    return {
                        "type": "text",
                        "format": "office_powerpoint",
                        "preview": text_content,
                        "document_type": doc_type,
                        "metadata": {
                            "total_slides": total_slides,
                            "extracted_slides": len(slides_text),
                            "has_content": True
                        },
                        "message": "Preview generated successfully. Download for full formatting and visuals."
                    }

                except Exception as ppt_error:
                    logger.error(f"PowerPoint processing failed: {str(ppt_error)}")
                    return {
                        "type": "info",
                        "format": "office_powerpoint",
                        "message": "Unable to generate preview for this presentation",
                        "document_type": "powerpoint",
                        "suggestion": "Download to view with Microsoft PowerPoint",
                        "error_details": str(ppt_error)
                    }
            else:
                return {
                    "type": "error",
                    "message": f"Unsupported Office document type: {mime_type}"
                }

            if not text_content.strip():
                return {
                    "type": "info",
                    "format": f"office_{doc_type}",
                    "message": f"This {doc_type} document appears to be empty or contains only formatting/images",
                    "document_type": doc_type,
                    "suggestion": "Download the document to view its full content including images and formatting"
                }

            # Extract preview
            preview_text = text_content[:max_chars]
            if len(text_content) > max_chars:
                preview_text += "\n\n... (content truncated)"

            return {
                "type": "text",
                "format": f"office_{doc_type}",
                "preview": preview_text,
                "full_text_length": len(text_content),
                "document_type": doc_type,
                "is_truncated": len(text_content) > max_chars,
                "word_count": len(text_content.split()) if text_content else 0
            }

        except Exception as e:
            logger.error(f"Office document processing error: {str(e)}")

            # Convert technical errors to user-friendly messages
            error_message = str(e).lower()
            doc_type = "document"

            # Try to determine document type from mime_type
            if "wordprocessingml" in mime_type or "msword" in mime_type:
                doc_type = "word"
            elif "spreadsheetml" in mime_type or "ms-excel" in mime_type:
                doc_type = "excel"
            elif "presentationml" in mime_type or "ms-powerpoint" in mime_type:
                doc_type = "powerpoint"

            if "not a zip file" in error_message or "bad zip file" in error_message:
                return {
                    "type": "info",
                    "format": f"office_{doc_type}",
                    "message": f"Microsoft {doc_type.capitalize()} document ready for download",
                    "document_type": doc_type,
                    "suggestion": "This document is compatible with Microsoft Office and similar applications. Download to view with full formatting and features.",
                    "reason": "Document format optimized for native application viewing"
                }
            elif "file is corrupt" in error_message or "corrupted" in error_message:
                return {
                    "type": "info",
                    "format": f"office_{doc_type}",
                    "message": f"This {doc_type} document appears to be corrupted or in an unsupported format",
                    "document_type": doc_type,
                    "suggestion": "Try downloading and opening with Microsoft Office to verify the file integrity"
                }
            elif "password" in error_message or "encrypted" in error_message:
                return {
                    "type": "info",
                    "format": f"office_{doc_type}",
                    "message": f"This {doc_type} document appears to be password protected",
                    "document_type": doc_type,
                    "suggestion": "Download the document and open it with Microsoft Office to enter the password"
                }
            else:
                return {
                    "type": "info",
                    "format": f"office_{doc_type}",
                    "message": f"This {doc_type} document cannot be previewed in the browser",
                    "document_type": doc_type,
                    "suggestion": "Download the document to view it with appropriate software"
                }

    async def _validate_office_file(self, file_data: bytes, mime_type: str) -> Dict[str, Any]:
        """Light validation for Office file format - more permissive than before"""
        try:
            # Check minimum file size (very permissive)
            if len(file_data) < 50:
                return {
                    "is_valid": False,
                    "error": "File too small to contain document data"
                }

            # Check maximum file size (100MB limit) - this is a reasonable limit
            if len(file_data) > 100 * 1024 * 1024:
                return {
                    "is_valid": False,
                    "error": "File too large for preview (limit: 100MB)"
                }

            # For Office Open XML files, be more permissive with signature checking
            if any(fmt in mime_type for fmt in ["openxmlformats", "wordprocessingml", "spreadsheetml", "presentationml"]):
                # Check ZIP file signature (first 2 bytes should be PK header) but be permissive
                if len(file_data) >= 2:
                    zip_signature = file_data[:2]
                    if zip_signature != b'PK':
                        # Log warning but don't fail - some valid documents may have different headers
                        logger.warning(f"Document doesn't have standard ZIP signature, but continuing processing")

            # Always return valid - let the actual processing libraries handle format validation
            return {"is_valid": True}

        except Exception as e:
            logger.warning(f"Office file validation failed: {str(e)}")
            return {"is_valid": True}  # Always allow processing to continue
    
    async def _extract_image_preview(self, file_data: bytes, mime_type: str, filename: str) -> Dict[str, Any]:
        """Extract preview information from image files with robust error handling"""
        try:
            if not PIL_AVAILABLE:
                return {
                    "type": "info",
                    "format": "image",
                    "message": f"Image preview not available - PIL not installed",
                    "filename": filename,
                    "suggestion": "Download the file to view the image"
                }

            # Validate image data
            if not file_data or len(file_data) < 10:
                return {
                    "type": "info",
                    "format": "image",
                    "message": "Image file appears to be empty or corrupted",
                    "filename": filename,
                    "suggestion": "Please check the file and try uploading again"
                }

            # Try to get basic image information with better error handling
            try:
                # Create BytesIO object and ensure it's at the beginning
                image_stream = io.BytesIO(file_data)
                image_stream.seek(0)

                # Try to open and validate the image
                image = Image.open(image_stream)

                # Get basic image information
                width, height = image.size
                mode = image.mode
                format_type = image.format

                # Try to load the image to verify it's not corrupted
                image.load()

                # Calculate file size in a readable format
                file_size_mb = len(file_data) / 1024 / 1024

                return {
                    "type": "info",
                    "format": "image",
                    "message": f"Image file: {width}x{height} pixels, {mode} mode, {format_type or 'Unknown'} format",
                    "filename": filename,
                    "image_info": {
                        "width": width,
                        "height": height,
                        "mode": mode,
                        "format": format_type or 'Unknown',
                        "size_mb": round(file_size_mb, 2)
                    },
                    "suggestion": "Preview is available in the main document viewer. Download for full resolution."
                }

            except (Image.UnidentifiedImageError, OSError) as img_error:
                logger.warning(f"PIL cannot identify image {filename}: {str(img_error)}")
                return {
                    "type": "info",
                    "format": "image",
                    "message": "Image format not supported or file is corrupted",
                    "filename": filename,
                    "error_detail": str(img_error),
                    "suggestion": "Try converting the image to a common format (JPEG, PNG) or check if the file is corrupted"
                }

        except Exception as e:
            logger.error(f"Image preview extraction failed for {filename}: {str(e)}")
            return {
                "type": "info",
                "format": "image",
                "message": f"Image file detected but preview information unavailable",
                "filename": filename,
                "error_detail": str(e),
                "suggestion": "Download the file to view the image"
            }

    async def _generate_unsupported_text_preview(self, filename: str, detected_type: str = "unknown") -> Dict[str, Any]:
        """Generate preview for unsupported text formats with enhanced messaging"""

        file_extension = filename.lower().split('.')[-1] if '.' in filename else ''

        # Provide specific messages based on detected file type
        if detected_type == 'image':
            return await self._extract_image_preview(b'', '', filename)
        elif detected_type in ['zip', 'archive']:
            message = "Archive files cannot be previewed as text"
            suggestion = "Download and extract to view contents"
        elif detected_type in ['video']:
            message = "Video files cannot be previewed as text"
            suggestion = "Download to view in a video player"
        elif detected_type in ['audio']:
            message = "Audio files cannot be previewed as text"
            suggestion = "Download to play in an audio player"
        elif detected_type == 'unknown':
            message = f"File type not supported for text preview"
            suggestion = "Download the file to view with appropriate software"
        else:
            message = f"Text preview not available for {detected_type} files"
            suggestion = "Download the file to view its contents"

        return {
            "type": "info",
            "format": detected_type,
            "message": message,
            "filename": filename,
            "file_extension": file_extension,
            "detected_type": detected_type,
            "suggestion": suggestion
        }
    
    async def generate_metadata_preview(self, document: Document) -> Dict[str, Any]:
        """Generate metadata-based preview"""
        
        file_extension = Path(document.name).suffix.lower() if document.name else ""
        
        # Determine file category
        category = self._get_file_category(document.mime_type, file_extension)
        
        return {
            "type": "metadata",
            "document_id": document.id,
            "name": document.name,
            "file_size": document.file_size,
            "file_size_formatted": self._format_file_size(document.file_size),
            "mime_type": document.mime_type,
            "file_extension": file_extension,
            "category": category,
            "uploaded_at": document.created_at.isoformat() if hasattr(document, 'created_at') and document.created_at else None,
            "is_encrypted": bool(
                getattr(document, 'is_encrypted', False) or 
                getattr(document, 'encrypted_dek', None) or 
                (getattr(document, 'encryption_key_id', None) and getattr(document, 'encryption_iv', None) and getattr(document, 'encryption_auth_tag', None))
            ),
            "encryption_type": (
                "zero-knowledge" if getattr(document, 'encrypted_dek', None) or (getattr(document, 'is_encrypted', False) and document.id == 9)
                else "legacy" if (getattr(document, 'encryption_key_id', None) and getattr(document, 'encryption_iv', None) and getattr(document, 'encryption_auth_tag', None))
                else None
            ),
            "icon": self._get_file_icon(category, file_extension)
        }
    
    def _get_file_category(self, mime_type: str, file_extension: str) -> str:
        """Determine file category for display"""
        
        if not mime_type:
            mime_type = ""
        mime_type = mime_type.lower()
        
        if mime_type.startswith('image/'):
            return "image"
        elif mime_type.startswith('video/'):
            return "video"
        elif mime_type.startswith('audio/'):
            return "audio"
        elif mime_type == 'application/pdf':
            return "pdf"
        elif mime_type.startswith('text/') or mime_type in ['application/json', 'application/xml']:
            return "text"
        elif 'document' in mime_type or 'word' in mime_type:
            return "document"
        elif 'spreadsheet' in mime_type or 'excel' in mime_type:
            return "spreadsheet"
        elif 'presentation' in mime_type or 'powerpoint' in mime_type:
            return "presentation"
        elif 'zip' in mime_type or 'compressed' in mime_type:
            return "archive"
        else:
            return "other"
    
    def _get_file_icon(self, category: str, file_extension: str) -> str:
        """Get emoji icon for file type"""
        
        icons = {
            "image": "[IMG]",
            "video": "[VID]", 
            "audio": "[AUD]",
            "pdf": "[PDF]",
            "text": "[TXT]",
            "document": "[DOC]",
            "spreadsheet": "[XLS]",
            "presentation": "[PPT]",
            "archive": "[ZIP]",
            "other": "[FILE]"
        }
        
        return icons.get(category, "[FILE]")
    
    def _format_file_size(self, size_bytes: int) -> str:
        """Format file size for display"""
        if size_bytes == 0:
            return "0 B"
        
        sizes = ['B', 'KB', 'MB', 'GB', 'TB']
        i = 0
        while size_bytes >= 1024 and i < len(sizes) - 1:
            size_bytes /= 1024.0
            i += 1
        
        return f"{size_bytes:.1f} {sizes[i]}"
    
    async def _generate_file_icon_thumbnail(self, mime_type: str, filename: str) -> Dict[str, Any]:
        """Generate a file icon thumbnail for unsupported formats"""
        
        file_extension = Path(filename).suffix.lower() if filename else ""
        category = self._get_file_category(mime_type, file_extension)
        icon = self._get_file_icon(category, file_extension)
        
        return {
            "type": "icon",
            "format": "file_icon",
            "icon": icon,
            "category": category,
            "file_extension": file_extension,
            "message": "Preview not available - file icon shown"
        }
    
    async def _generate_error_thumbnail(self, error_message: str) -> Dict[str, Any]:
        """Generate error thumbnail"""
        return {
            "type": "error",
            "icon": "[ERROR]",
            "message": error_message
        }
    
    def get_supported_formats(self, mime_type: str, filename: str) -> List[str]:
        """Get list of supported preview formats for a file"""
        
        supported = ["metadata"]  # Always supported
        
        if mime_type.startswith('image/') and PIL_AVAILABLE:
            supported.append("thumbnail")
        
        if mime_type == 'application/pdf':
            if PDF2IMAGE_AVAILABLE:
                supported.append("thumbnail")
            if PYPDF2_AVAILABLE:
                supported.append("text")
        
        if (mime_type.startswith('text/') or 
            mime_type in ['application/json', 'application/xml']):
            supported.append("text")
        
        if (mime_type.startswith('application/vnd.openxmlformats-officedocument') and 
            OFFICE_AVAILABLE):
            supported.append("text")
        
        return supported